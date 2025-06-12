"""
AI Service for document analysis and summarization.
Handles text extraction from various file formats and Google Gemini integration.
"""

import google.generativeai as genai
from django.conf import settings
from docx import Document
import openpyxl
import io
import os
from typing import Optional


class AIService:
    """Service class for AI-powered document analysis"""
    
    def __init__(self):
        """Initialize Gemini client with API key from settings"""
        if not settings.GEMINI_API_KEY or settings.GEMINI_API_KEY == 'your-gemini-api-key-here':
            # Demo mode - no real API key configured
            self.demo_mode = True
            self.model = None
            print("AI Service running in demo mode - no real Gemini API key configured")
        else:
            self.demo_mode = False
            try:
                genai.configure(api_key=settings.GEMINI_API_KEY)
                self.model = genai.GenerativeModel('gemini-1.5-flash')
                print("AI Service initialized with Google Gemini API key")
            except Exception as e:
                print(f"Failed to initialize Gemini client: {e}")
                print("Falling back to demo mode")
                self.demo_mode = True
                self.model = None

    def extract_text_from_file(self, file_path: str) -> str:
        """
        Extract text content from uploaded file based on file extension.
        
        Args:
            file_path: Path to the file to extract text from
            
        Returns:
            Extracted text content as string
            
        Raises:
            ValueError: If file type is not supported
            Exception: If file reading fails
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        file_extension = os.path.splitext(file_path)[1].lower()
        
        try:
            if file_extension == '.docx':
                return self._extract_from_docx(file_path)
            elif file_extension == '.xlsx':
                return self._extract_from_xlsx(file_path)
            elif file_extension == '.pptx':
                return self._extract_from_pptx(file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_extension}")
        except Exception as e:
            raise Exception(f"Failed to extract text from {file_extension} file: {str(e)}")

    def _extract_from_docx(self, file_path: str) -> str:
        """Extract text from DOCX file"""
        doc = Document(file_path)
        text_content = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content.append(paragraph.text.strip())
        
        return '\n'.join(text_content)

    def _extract_from_xlsx(self, file_path: str) -> str:
        """Extract text from XLSX file"""
        workbook = openpyxl.load_workbook(file_path, data_only=True)
        text_content = []
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text_content.append(f"Sheet: {sheet_name}")
            
            for row in sheet.iter_rows(values_only=True):
                row_text = []
                for cell in row:
                    if cell is not None:
                        row_text.append(str(cell))
                if row_text:
                    text_content.append(' | '.join(row_text))
        
        workbook.close()
        return '\n'.join(text_content)

    def _extract_from_pptx(self, file_path: str) -> str:
        """Extract text from PPTX file"""
        # For PPTX, we'll use python-pptx library
        from pptx import Presentation
        
        presentation = Presentation(file_path)
        text_content = []
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            text_content.append(f"Slide {slide_num}:")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text.strip())
        
        return '\n'.join(text_content)

    def _truncate_text(self, text: str) -> str:
        """
        Truncate text to fit within Gemini token limits.
        Gemini has a generous context window, but we'll still limit for performance
        """
        max_chars = 30000  # Conservative limit for Gemini
        if len(text) <= max_chars:
            return text
        
        # Truncate and add notice
        truncated = text[:max_chars]
        return f"{truncated}\n\n[Note: Content truncated due to length limits]"

    def summarize_document(self, text: str) -> str:
        """
        Generate AI summary of document content using Google Gemini.
        
        Args:
            text: Extracted text content from document
            
        Returns:
            AI-generated summary
            
        Raises:
            Exception: If Gemini API call fails
        """
        if not text.strip():
            return "No text content found in the document."
        
        # Check if in demo mode
        if self.demo_mode:
            # Return a demo summary based on content analysis
            word_count = len(text.split())
            char_count = len(text)
            lines = text.split('\n')
            non_empty_lines = [line for line in lines if line.strip()]
            
            return f"""**Demo AI Summary**

This document contains approximately {word_count} words and {char_count} characters across {len(non_empty_lines)} lines of content.

**Content Analysis:**
- Document appears to be a {'spreadsheet' if 'Sheet:' in text else 'text document'}
- Contains structured data with multiple entries
- Includes various fields and data points

**Key Observations:**
- The document has tabular data organization
- Multiple columns of information are present
- Contains what appears to be student or record management data

*Note: This is a demo summary. Configure GEMINI_API_KEY for full AI analysis.*"""
        
        # Truncate text if necessary
        truncated_text = self._truncate_text(text)
        
        prompt = f"""Please provide a clear and concise summary of this document. Focus on:
1. The document's main purpose and content type
2. Key information and data points
3. Structure and organization
4. Any important patterns or insights

Document content:
{truncated_text}"""
        
        try:
            response = self.model.generate_content(prompt)
            return response.text
            
        except Exception as e:
            # If it's a quota/rate limit error, provide a helpful message
            if "quota" in str(e).lower() or "rate limit" in str(e).lower():
                return f"""**AI Summary (Rate Limited)**

This document contains approximately {len(text.split())} words and appears to be a {'spreadsheet' if 'Sheet:' in text else 'text document'}.

**Note**: Gemini API rate limit exceeded. Please try again in a moment.

**Demo Analysis**: The document appears to contain structured data with multiple fields and entries.

*For full AI analysis, please wait for the rate limit to reset.*"""
            else:
                raise Exception(f"AI summarization failed: {str(e)}")

    def generate_file_summary(self, file_path: str) -> str:
        """
        Complete workflow: extract text and generate summary.
        
        Args:
            file_path: Path to the file to summarize
            
        Returns:
            AI-generated summary of the file
        """
        try:
            # Extract text from file
            text = self.extract_text_from_file(file_path)
            
            # Generate summary using AI
            summary = self.summarize_document(text)
            
            return summary
            
        except Exception as e:
            raise Exception(str(e))
